#!/usr/bin/env python3
"""Lab 1 — Cloud Compute intro deck. Ocean Gradient design system, ported to python-pptx."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- palette ----------
NAVY="0B2942"; MIDNIGHT="21295C"; DEEPBLUE="065A82"; TEAL="1C7293"
ICE="CFE8F0"; WHITE="FFFFFF"; INK="1B2733"; MUTE="5C7080"; CARD="F2F8FA"
DARKCARD="16385A"; NAVYCARD="0F2036"
TITLEF="Cambria"; BODYF="Calibri"; MONOF="Courier New"

def C(h): return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))

prs=Presentation()
prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]

def slide(bg=WHITE):
    s=prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb=C(bg)
    return s

def _set_runs(tf, runs, size, color, bold, font, italic, align, line_spacing):
    p=tf.paragraphs[0]
    if align is not None: p.alignment=align
    if line_spacing: p.line_spacing=line_spacing
    if isinstance(runs,str): runs=[(runs,{})]
    for text,opt in runs:
        r=p.add_run(); r.text=text
        r.font.size=Pt(opt.get("size",size)); r.font.bold=opt.get("bold",bold)
        r.font.italic=opt.get("italic",italic); r.font.name=opt.get("font",font)
        r.font.color.rgb=C(opt.get("color",color))
        if opt.get("link"): r.hyperlink.address=opt["link"]

def text(s, runs, x,y,w,h, size=12, color=INK, bold=False, font=BODYF,
         italic=False, align=None, anchor=MSO_ANCHOR.TOP, line_spacing=None, wrap=True):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=wrap; tf.vertical_anchor=anchor
    for m in ("margin_left","margin_right","margin_top","margin_bottom"): setattr(tf,m,0)
    _set_runs(tf,runs,size,color,bold,font,italic,align,line_spacing)
    return tb

def _noline(sh): sh.line.fill.background()
def _noshadow(sh):
    try: sh.shadow.inherit=False
    except Exception: pass

def rrect(s, x,y,w,h, fill, line=None, lw=0.75, radius=0.08):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    try: sh.adjustments[0]=radius
    except Exception: pass
    sh.fill.solid(); sh.fill.fore_color.rgb=C(fill)
    if line: sh.line.color.rgb=C(line); sh.line.width=Pt(lw)
    else: _noline(sh)
    _noshadow(sh); return sh

def rect(s, x,y,w,h, fill, line=None, lw=0.75):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=C(fill)
    if line: sh.line.color.rgb=C(line); sh.line.width=Pt(lw)
    else: _noline(sh)
    _noshadow(sh); return sh

def oval(s, x,y,w,h, fill, line=None, lw=0.75):
    sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=C(fill)
    if line: sh.line.color.rgb=C(line); sh.line.width=Pt(lw)
    else: _noline(sh)
    _noshadow(sh); return sh

def hline(s, x,y,w, color, lw=1.5):
    ln=s.shapes.add_connector(2,Inches(x),Inches(y),Inches(x+w),Inches(y))
    ln.line.color.rgb=C(color); ln.line.width=Pt(lw); return ln

def chip(s, label, x,y,size, bg, fg=WHITE):
    """colored circle with a white letter/number — always high-contrast."""
    oval(s,x,y,size,size,bg)
    text(s,label,x,y-0.02,size,size,size=int(size*20),color=fg,bold=True,
         font=TITLEF,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

def kicker(s, t, color=TEAL):
    text(s,t.upper(),0.6,0.42,10,0.35,size=13,color=color,bold=True,font=BODYF)

def title(s, t, color=INK, size=32, w=11.7):
    text(s,t,0.6,0.76,w,0.95,size=size,color=color,bold=True,font=TITLEF)

def pathtag(s, t, x,y,w, dark=False):
    rrect(s,x,y,w,0.34,DARKCARD if dark else ICE,("3F6E8C" if dark else "AEC7D6"),0.75,0.14)
    text(s,t,x+0.12,y,w-0.24,0.34,size=10.5,color=("8FD4E8" if dark else DEEPBLUE),
         bold=True,font=MONOF,anchor=MSO_ANCHOR.MIDDLE)

def backtag(s, t, x,y,w, dark=False):
    rrect(s,x,y,w,0.34,("2A3F63" if dark else "E9F1E6"),("5A8FB0" if dark else "8FB88F"),0.75,0.24)
    text(s,"↩  "+t,x+0.16,y,w-0.3,0.34,size=10,italic=True,
         color=("BFDCE8" if dark else "3D6B3D"),font=BODYF,anchor=MSO_ANCHOR.MIDDLE)

def footer(s, dark=False, n=None):
    text(s,"LAB 1 · CLOUD COMPUTE · AI SEMINAR",0.6,7.06,7,0.3,
         size=9,color=("8FA8C2" if dark else MUTE),font=BODYF)
    if n is not None:
        text(s,str(n),12.5,7.03,0.6,0.3,size=10,color=("8FA8C2" if dark else MUTE),
             font=BODYF,align=PP_ALIGN.RIGHT)

# ===================== SLIDE 1: TITLE =====================
s=slide(NAVY)
oval(s,10.4,-1.5,5.2,5.2,MIDNIGHT); oval(s,12.1,4.4,3.4,3.4,DEEPBLUE)
# small node/edge decoration
nodes=[(10.7,1.5),(11.6,1.05),(11.6,2.0),(12.5,0.7),(12.5,1.4),(12.5,1.95),(12.5,2.5)]
edges=[(0,1),(0,2),(1,3),(1,4),(2,5),(2,6)]
for a,b in edges:
    x1,y1=nodes[a]; x2,y2=nodes[b]
    ln=s.shapes.add_connector(2,Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    ln.line.color.rgb=C("5A8FB0"); ln.line.width=Pt(1.25)
for x,y in nodes: oval(s,x-0.06,y-0.06,0.12,0.12,"8FD4E8")
text(s,"AI SEMINAR  ·  LAB 1",0.8,1.5,8,0.4,size=14,color="7FB8D6",bold=True,font=BODYF)
text(s,"Cloud Compute",0.75,2.0,11,1.4,size=60,color=WHITE,bold=True,font=TITLEF)
text(s,"Running Models You Can't Run on Your Laptop — Kaggle & RunPod",
     0.8,3.35,10.5,0.7,size=20,color=ICE,font=BODYF)
hline(s,0.8,4.28,3.2,TEAL,2)
text(s,[("How we'll get there:  ",{"bold":True,"color":"BFE0EE"}),
        ("why this is a lab and not a lesson, what cloud compute for ML actually is, a one-slide tour of each platform, a side-by-side of when to use which, then a preview of the two hands-on parts you'll work through.",{"color":"9FC3D9"})],
     0.8,4.5,8.7,1.3,size=13.5,font=BODYF,line_spacing=1.25)
text(s,"Companion to Lesson 3: repeng   ·   github.com/vgel/repeng",
     0.8,6.6,9,0.4,size=12,color="6FA0BE",font=MONOF)

# ===================== SLIDE 2: WHY A LAB =====================
s=slide(WHITE)
kicker(s,"Orientation"); title(s,"This Is a Lab, Not a Lesson")
text(s,[("Lesson 2 (LogitLoom) ran entirely in your browser against a hosted API — no hardware required. ",{"color":MUTE}),
        ("repeng is different: it downloads and runs the actual model weights. ",{"color":INK,"bold":True}),
        ("Mistral-7B",{"color":DEEPBLUE,"bold":True,"link":"https://huggingface.co/mistralai/Mistral-7B-Instruct-v0.1"}),
        (" is ~7.24B parameters — about 14.5 GB in fp16 — needing far more GPU memory (VRAM) than the 4–8 GB on a typical laptop GPU, and many laptops have no dedicated GPU at all.",{"color":MUTE})],
     0.6,1.5,12.2,1.1,size=13.5,font=BODYF,line_spacing=1.22)
# two definition cards
rrect(s,0.6,2.6,5.95,3.15,CARD,"D9E6EC",0.75)
chip(s,"L",0.9,2.9,0.7,DEEPBLUE)
text(s,"Lesson",1.75,2.95,4.5,0.6,size=18,bold=True,color=INK,font=BODYF,anchor=MSO_ANCHOR.MIDDLE)
text(s,"Teaches one repo / tool in depth — what it is, how it works, how to use it. LogitLoom and repeng are lessons.",
     0.95,3.8,5.25,1.8,size=13,color=MUTE,font=BODYF,line_spacing=1.25)
rrect(s,6.75,2.6,5.95,3.15,MIDNIGHT)
chip(s,"L",7.05,2.9,0.7,TEAL)
text(s,"Lab",7.9,2.95,4.5,0.6,size=18,bold=True,color=WHITE,font=BODYF,anchor=MSO_ANCHOR.MIDDLE)
text(s,"Teaches a foundational platform or skill you need before you can explore certain repos. Cloud compute is Lab 1.",
     7.1,3.8,5.25,1.8,size=13,color="D7E9F2",font=BODYF,line_spacing=1.25)
text(s,[("The point: ",{"bold":True,"color":INK}),
        ("before the repeng lesson, you need somewhere to run models. This lab builds that muscle — in two parts, Kaggle (A) and RunPod (B).",{"color":MUTE})],
     0.6,6.05,12.1,0.7,size=13,font=BODYF,line_spacing=1.2)
footer(s,n=2)

# ===================== SLIDE 3: LANDSCAPE =====================
s=slide(WHITE)
kicker(s,"Background"); title(s,"Cost Is the Only Sharp Divide")
text(s,"Three dimensions to weigh. On cost, Kaggle and RunPod are genuinely opposite; on the other two they lean different ways but overlap.",
     0.6,1.55,12.2,0.5,size=14,color=MUTE,font=BODYF)
axes=[("1","Ephemeral vs. persistent","A hosted notebook that spins up, times out, and resets — versus a machine you rent and keep alive while you need it."),
      ("2","Free vs. paid","A quota-limited free tier on shared hardware — versus pay-per-use where the GPU (and the clock) is yours."),
      ("3","Managed vs. DIY","An environment handed to you ready to go — versus one you install, configure, and manage yourself.")]
cw=3.9; gap=0.22; x0=0.6; y=2.35; ch=3.5
for i,(num,h,b) in enumerate(axes):
    x=x0+i*(cw+gap)
    rrect(s,x,y,cw,ch,CARD,"D9E6EC",0.75)
    chip(s,num,x+cw/2-0.4,y+0.35,0.8,DEEPBLUE)
    text(s,h,x+0.25,y+1.35,cw-0.5,0.6,size=15.5,bold=True,color=INK,font=BODYF,align=PP_ALIGN.CENTER)
    text(s,b,x+0.3,y+2.0,cw-0.6,1.3,size=12,color=MUTE,font=BODYF,align=PP_ALIGN.CENTER,line_spacing=1.2)
text(s,[("Bottom line: ",{"bold":True,"color":INK}),
        ("the real either/or is cost — Kaggle free (with limits), RunPod pay-as-you-go. On persistence and how hands-on you get, both cover a range and overlap (both run notebooks and scripts). Either way: a GPU your laptop doesn't have.",{"color":MUTE})],
     0.6,6.05,12.2,0.8,size=13,font=BODYF,line_spacing=1.2)
footer(s,n=3)

# ===================== SLIDE 4: KAGGLE =====================
s=slide(WHITE)
kicker(s,"Platform A"); title(s,"Kaggle")
chip(s,"K",11.75,0.6,0.85,DEEPBLUE)
text(s,"A free, browser-based notebook platform (by Google) with datasets, pre-hosted models, and a big data-science community.",
     0.6,1.55,11.0,0.6,size=14,color=MUTE,font=BODYF,line_spacing=1.2)
# what you get / limits, two columns
rrect(s,0.6,2.35,5.95,4.2,CARD,"D9E6EC",0.75)
text(s,"What you get",0.9,2.55,5.3,0.4,size=15,bold=True,color=DEEPBLUE,font=BODYF)
gets=[("Free GPUs","T4×2 or P100, within a weekly quota — what we'll use."),
      ("(TPUs too)","also available on Kaggle — we'll come back to those in a later lesson."),
      ("Notebooks","Fork, run, and version Jupyter-style notebooks in the browser."),
      ("Datasets","Upload any files and mount them read-only — how we'll ship the repeng wheels."),
      ("Models","Pre-hosted weights (incl. Mistral) added from a sidebar search."),
      ("Secrets","Store tokens (e.g. an HF token) without hard-coding them.")]
gy=2.92
for h,b in gets:
    text(s,[(h+":  ",{"bold":True,"color":INK}),(b,{"color":MUTE})],
         0.95,gy,5.3,0.58,size=11.5,font=BODYF,line_spacing=1.1); gy+=0.6
rrect(s,6.75,2.35,5.95,4.2,MIDNIGHT)
text(s,"Limits to plan around",7.05,2.55,5.3,0.4,size=15,bold=True,color="8FD4E8",font=BODYF)
lims=[("Sessions time out","after a fixed window — long runs can get cut off."),
      ("Not persistent","the environment resets between sessions; save what you need."),
      ("No root","you can pip-install, but not reconfigure the machine."),
      ("Internet off by default","toggle it on, or install offline from a Dataset (what Part A does).")]
ly=3.0
for h,b in lims:
    text(s,[(h+"  ",{"bold":True,"color":"D7E9F2"}),("— "+b,{"color":"AECBDD"})],
         7.1,ly,5.3,0.75,size=11.5,font=BODYF,line_spacing=1.15); ly+=0.85
footer(s,n=4)

# ===================== SLIDE 5: RUNPOD =====================
s=slide(NAVY)
kicker(s,"Platform B","8FD4E8"); title(s,"RunPod",WHITE)
chip(s,"R",11.75,0.6,0.85,TEAL)
text(s,"Rent a GPU “pod” — a container on a cloud GPU — by the second or hour, and connect to it via Jupyter or SSH.",
     0.6,1.55,11.0,0.6,size=14,color="9FC3D9",font=BODYF,line_spacing=1.2)
rrect(s,0.6,2.35,5.95,4.2,MIDNIGHT)
text(s,"What you get",0.9,2.55,5.3,0.4,size=15,bold=True,color="8FD4E8",font=BODYF)
gets=[("Your choice of GPU","from small cards up to A100 / H100."),
      ("Jupyter or a terminal","run notebooks, or plain scripts — your call."),
      ("Persistent storage","volumes that survive while the pod lives."),
      ("Full control","install anything; it's your environment."),
      ("Real repos","clone from GitHub, work, and push back.")]
gy=3.0
for h,b in gets:
    text(s,[(h+":  ",{"bold":True,"color":"D7E9F2"}),(b,{"color":"AECBDD"})],
         0.95,gy,5.3,0.7,size=11.5,font=BODYF,line_spacing=1.12); gy+=0.72
rrect(s,6.75,2.35,5.95,4.2,"16324E")
text(s,"Limits to plan around",7.05,2.55,5.3,0.4,size=15,bold=True,color="FFB4A2",font=BODYF)
lims=[("It costs money","you pay while the pod runs — remember to stop it."),
      ("You manage setup","dependencies and drivers are on you."),
      ("Persistence needs care","use a volume, or your files vanish when the pod dies."),
      ("Getting files out","you have to pull outputs off the box deliberately.")]
ly=3.0
for h,b in lims:
    text(s,[(h+"  ",{"bold":True,"color":"D7E9F2"}),("— "+b,{"color":"AECBDD"})],
         7.1,ly,5.3,0.75,size=11.5,font=BODYF,line_spacing=1.15); ly+=0.85
footer(s,dark=True,n=5)

# ===================== SLIDE 6: WHEN TO USE WHICH (centerpiece) =====================
s=slide(WHITE)
kicker(s,"The Decision"); title(s,"Match the Platform to the Job")
text(s,"Same job, two very different fits. Screenshot this one.",
     0.6,1.5,12,0.4,size=13.5,color=MUTE,font=BODYF)
rows=[("","Kaggle","RunPod"),
      ("Cost","Free (weekly quota)","Pay per second / hour"),
      ("Best for","Quick experiments, learning, sharing","Big models, long runs, full control"),
      ("Hardware","Fixed free GPUs (T4 / P100)","Pick your GPU, up to A100 / H100"),
      ("Environment","Managed; resets each session","Yours to set up; persists with the pod"),
      ("Internet","Off by default","On"),
      ("Workflow","Notebooks (fork / import)","Notebooks or scripts, your own repo")]
tx,ty,tw=0.6,2.05,12.13
rowh=[0.5]+[0.6]*6
th=sum(rowh)
gtbl=s.shapes.add_table(len(rows),3,Inches(tx),Inches(ty),Inches(tw),Inches(th)).table
gtbl.first_row=False; gtbl.horz_banding=False
colw=[3.0,4.56,4.57]
for j,wd in enumerate(colw): gtbl.columns[j].width=Inches(wd)
for i in range(len(rows)): gtbl.rows[i].height=Inches(rowh[i])
def style_cell(cell, txt, size, color, bold, fill, align=PP_ALIGN.LEFT):
    cell.fill.solid(); cell.fill.fore_color.rgb=C(fill)
    cell.vertical_anchor=MSO_ANCHOR.MIDDLE
    cell.margin_left=Inches(0.15); cell.margin_right=Inches(0.1)
    cell.margin_top=Inches(0.02); cell.margin_bottom=Inches(0.02)
    tf=cell.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=txt; r.font.size=Pt(size); r.font.bold=bold
    r.font.name=BODYF; r.font.color.rgb=C(color)
for i,row in enumerate(rows):
    for j,val in enumerate(row):
        if i==0:
            fill=[CARD,DEEPBLUE,TEAL][j]
            col=[MUTE,WHITE,WHITE][j]
            style_cell(gtbl.cell(i,j),val,15,col,True,fill,
                       PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER)
        else:
            if j==0:
                style_cell(gtbl.cell(i,j),val,12.5,INK,True,"E9F1F5")
            else:
                fill=WHITE if i%2 else CARD
                style_cell(gtbl.cell(i,j),val,12,MUTE,False,fill,PP_ALIGN.CENTER)
# rule-of-thumb banner
rrect(s,0.6,6.35,12.13,0.62,MIDNIGHT)
text(s,[("Rule of thumb:  ",{"bold":True,"color":"8FD4E8"}),
        ("prototyping, free, shareable → Kaggle.   Heavy, long, or you need control → RunPod.",{"color":"E6F1F6"})],
     0.85,6.35,11.6,0.62,size=13,font=BODYF,anchor=MSO_ANCHOR.MIDDLE)
footer(s,n=6)

# ===================== SLIDE 7: PART A PREVIEW =====================
s=slide(WHITE)
kicker(s,"Part A · Kaggle"); title(s,"Four Steps to a Saved Control Vector")
beats=[("1","Add a model","Pull Mistral-7B from Kaggle Models — hosted, license handled in-platform, just works."),
       ("2","Install repeng offline","From an uploaded wheels Dataset, with the internet turned off."),
       ("3","Hit & fix a token error","Pull the same model from Hugging Face, get the gated-repo error, fix it with Kaggle Secrets."),
       ("4","Train & save","Build a control vector with repeng, export it, and save your notebook output.")]
cw=2.95; gap=0.19; x0=0.6; y=2.05; ch=3.7
for i,(num,h,b) in enumerate(beats):
    x=x0+i*(cw+gap)
    rrect(s,x,y,cw,ch,CARD,"D9E6EC",0.75)
    chip(s,num,x+cw/2-0.38,y+0.35,0.76,DEEPBLUE)
    text(s,h,x+0.22,y+1.3,cw-0.44,0.7,size=15,bold=True,color=INK,font=BODYF,align=PP_ALIGN.CENTER)
    text(s,b,x+0.26,y+2.05,cw-0.52,1.5,size=11.5,color=MUTE,font=BODYF,align=PP_ALIGN.CENTER,line_spacing=1.18)
text(s,[("On purpose: ",{"bold":True,"color":INK}),
        ("the code is provided and the ML is trivial — the skill you're building is operating Kaggle, not the modeling.",{"color":MUTE})],
     0.6,6.1,12.1,0.7,size=13,font=BODYF,line_spacing=1.2)
footer(s,n=7)

# ===================== SLIDE 8: PART B PREVIEW =====================
s=slide(NAVY)
kicker(s,"Part B · RunPod","8FD4E8"); title(s,"Clone, Break, Fix, Push",WHITE)
beats=[("1","Launch a pod","Pick a GPU, start it, and open Jupyter."),
       ("2","Clone a repo","Fork it on GitHub, then git clone onto the pod."),
       ("3","Fix what breaks","Run a script that works locally but breaks on RunPod — twice — and recover each time."),
       ("4","Push & pull","Commit your change back to your fork, and pull your output files off the box.")]
cw=2.95; gap=0.19; x0=0.6; y=2.05; ch=3.7
for i,(num,h,b) in enumerate(beats):
    x=x0+i*(cw+gap)
    rrect(s,x,y,cw,ch,MIDNIGHT)
    chip(s,num,x+cw/2-0.38,y+0.35,0.76,TEAL)
    text(s,h,x+0.22,y+1.3,cw-0.44,0.7,size=15,bold=True,color=WHITE,font=BODYF,align=PP_ALIGN.CENTER)
    text(s,b,x+0.26,y+2.05,cw-0.52,1.5,size=11.5,color="AECBDD",font=BODYF,align=PP_ALIGN.CENTER,line_spacing=1.18)
text(s,[("Don't forget: ",{"bold":True,"color":"8FD4E8"}),
        ("the meter runs while a pod is alive. Stopping it when you're done is part of the lesson.",{"color":"AECBDD"})],
     0.6,6.1,12.1,0.7,size=13,font=BODYF,line_spacing=1.2)
footer(s,dark=True,n=8)

# ===================== SLIDE 9: GROUND RULES =====================
s=slide(WHITE)
kicker(s,"Before You Start"); title(s,"Set Up Your Accounts First")
# accounts row
text(s,"Accounts to create",0.6,1.55,6,0.4,size=15,bold=True,color=DEEPBLUE,font=BODYF)
accts=[("Kaggle","free; verify your phone to unlock GPU + internet"),
       ("RunPod","needs a little prepaid credit"),
       ("GitHub","for the Part B clone / push round-trip"),
       ("Hugging Face","for a token + accepting a model license")]
ay=2.05
for h,b in accts:
    oval(s,0.65,ay+0.06,0.13,0.13,TEAL)
    text(s,[(h+"  ",{"bold":True,"color":INK}),("— "+b,{"color":MUTE})],
         0.95,ay-0.05,5.6,0.5,size=12.5,font=BODYF,line_spacing=1.15); ay+=0.62
# HF license callout card
rrect(s,6.75,1.9,5.95,2.7,CARD,"D9E6EC",0.75)
chip(s,"!",7.05,2.15,0.6,"B85042")
text(s,"About that model license",7.8,2.18,4.7,0.55,size=15,bold=True,color=INK,font=BODYF,anchor=MSO_ANCHOR.MIDDLE)
text(s,"Before Hugging Face lets you download some models (Mistral included), you have to accept a license. That's not a formality to click past — it's a real gate on real model weights, and knowing it exists is part of the point.",
     7.1,2.95,5.3,1.5,size=12,color=MUTE,font=BODYF,line_spacing=1.22)
# where the files are
rrect(s,0.6,4.95,12.13,1.05,MIDNIGHT)
text(s,[("Where the files live:  ",{"bold":True,"color":"8FD4E8"}),
        ("lab-1-cloud-compute/  — start with Part A (kaggle_assignment.md), then Part B (runpod_assignment.md).",{"color":"E6F1F6","font":MONOF})],
     0.9,5.15,11.5,0.65,size=12.5,font=BODYF,anchor=MSO_ANCHOR.MIDDLE,line_spacing=1.15)
text(s,"Goal: by the end, running a model on borrowed hardware feels routine — so the repeng lesson can be about repeng.",
     0.6,6.25,12.1,0.6,size=13,italic=True,color=MUTE,font=BODYF,line_spacing=1.2)
footer(s,dark=True,n=9)

out="/sessions/nifty-exciting-davinci/mnt/outputs/Lab1_CloudCompute_Intro.pptx"
prs.save(out)
print("SAVED",out,"slides:",len(prs.slides._sldIdLst))
